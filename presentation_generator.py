"""
Presentation Generator

This module creates professional PowerPoint presentations from markdown content.
It intelligently analyzes content, distributes it across slides appropriately,
and applies professional design elements.
"""

import os
import re
import glob
from pathlib import Path
import argparse
import markdown
import nltk
from nltk.tokenize import sent_tokenize
from bs4 import BeautifulSoup
import colorsys
import random
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Try to download NLTK data if needed
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# Constants for presentation design
SLIDE_WIDTH = Inches(16)
SLIDE_HEIGHT = Inches(9)  # 16:9 aspect ratio

COLOR_SCHEMES = {
    "professional": {
        "primary": RGBColor(0, 76, 153),
        "secondary": RGBColor(102, 153, 255),
        "accent": RGBColor(255, 153, 51),
        "background": RGBColor(255, 255, 255),
        "text": RGBColor(0, 0, 0),
        "light_text": RGBColor(100, 100, 100)
    },
    "modern": {
        "primary": RGBColor(66, 133, 244),
        "secondary": RGBColor(234, 67, 53),
        "accent": RGBColor(251, 188, 5),
        "background": RGBColor(255, 255, 255),
        "text": RGBColor(32, 33, 36),
        "light_text": RGBColor(95, 99, 104)
    },
    "minimal": {
        "primary": RGBColor(0, 0, 0),
        "secondary": RGBColor(128, 128, 128),
        "accent": RGBColor(212, 33, 40),
        "background": RGBColor(250, 250, 250),
        "text": RGBColor(50, 50, 50),
        "light_text": RGBColor(150, 150, 150)
    }
}

# Slide content limits for optimal readability
MAX_BULLET_POINTS_PER_SLIDE = 6
MAX_WORDS_PER_SLIDE = 120
MAX_CHARS_PER_BULLET = 100
MAX_TABLE_ROWS_PER_SLIDE = 8

class PresentationGenerator:
    def __init__(self, company_dir, color_scheme="professional", theme="modern"):
        """
        Initialize the presentation generator.
        
        Args:
            company_dir (str): Directory containing markdown files
            color_scheme (str): Color scheme name from predefined schemes
            theme (str): Design theme name
        """
        self.company_dir = company_dir
        self.color_scheme = COLOR_SCHEMES[color_scheme]
        self.theme = theme
        
        # Extract company name from directory
        self.company_name = os.path.basename(company_dir).split('_')[0]
        
        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        
        # Create slide layouts
        self._create_slide_layouts()
        
        # Sections and their content
        self.sections = {}
        
    def _create_slide_layouts(self):
        """Create slide layouts for the presentation."""
        # Store slide layout references for convenience
        layouts = self.prs.slide_layouts
        self.layouts = {
            'title': layouts[0],  # Title slide layout
            'title_content': layouts[1],  # Title and content layout
            'section': layouts[2],  # Section header layout
            'two_content': layouts[3],  # Two content layout
            'title_only': layouts[5],  # Title only layout
            'blank': layouts[6]  # Blank layout
        }
        
    def _extract_content_from_markdown(self):
        """Extract content from markdown files in the company directory."""
        # Look for markdown files in the markdown directory
        markdown_dir = os.path.join(self.company_dir, "markdown")
        if not os.path.exists(markdown_dir):
            raise FileNotFoundError(f"Markdown directory not found: {markdown_dir}")
        
        markdown_files = glob.glob(os.path.join(markdown_dir, "*.md"))
        
        # Import config.py for section order and source patterns
        import sys
        sys.path.append(os.path.dirname(os.path.dirname(self.company_dir)))
        
        try:
            from config import SECTION_ORDER, PDF_CONFIG
            self.section_order = [item[0] for item in SECTION_ORDER]
            self.section_titles = {item[0]: item[1] for item in SECTION_ORDER}
            self.source_patterns = PDF_CONFIG['SOURCES']['SOURCE_HEADING_PATTERNS']
        except ImportError:
            # Default section order if config.py can't be imported
            self.section_order = [
                "basic", "vision", "management_strategy", "management_message", 
                "crisis", "digital_transformation", "financial", "competitive", 
                "regulatory", "business_structure", "strategy_research"
            ]
            self.section_titles = {
                "basic": "Basic Information",
                "vision": "Vision Analysis",
                "management_strategy": "Management Strategy",
                "management_message": "Management Message",
                "crisis": "Crisis Management",
                "digital_transformation": "Digital Transformation Analysis",
                "financial": "Financial Analysis",
                "competitive": "Competitive Landscape",
                "regulatory": "Regulatory Environment",
                "business_structure": "Business Structure",
                "strategy_research": "Strategy Research"
            }
            self.source_patterns = ["Sources", "References", "Bibliography"]
        
        # Process each markdown file
        for md_file in markdown_files:
            section_name = os.path.basename(md_file).replace('.md', '')
            if section_name in self.section_order:
                with open(md_file, 'r', encoding='utf-8') as f:
                    content = f.read()
                    
                # Remove sources sections
                content = self._remove_sources(content)
                
                # Convert markdown to HTML for easier parsing
                html_content = markdown.markdown(content, extensions=['tables', 'fenced_code'])
                
                # Store processed content
                self.sections[section_name] = {
                    'title': self.section_titles.get(section_name, section_name.replace('_', ' ').title()),
                    'html': html_content,
                    'raw': content
                }
    
    def _remove_sources(self, content):
        """Remove sources sections from markdown content."""
        result = content
        
        # Try to find and remove source sections
        for pattern in self.source_patterns:
            # Look for headings with different levels (# to ####)
            for heading_level in range(1, 5):
                heading = '#' * heading_level + ' ' + pattern
                parts = result.split(heading)
                if len(parts) > 1:
                    # Keep everything before the sources section
                    # Find the next heading of same or higher level
                    next_heading_match = re.search(r'^#{1,' + str(heading_level) + r'}\s', parts[1], re.MULTILINE)
                    if next_heading_match:
                        # Keep content up to the next heading
                        parts[1] = parts[1][next_heading_match.start():]
                        result = parts[0] + parts[1]
                    else:
                        # No next heading, just take the first part
                        result = parts[0]
        
        return result
    
    def _analyze_content(self, html_content):
        """
        Analyze HTML content to determine how to split it into slides.
        
        Returns:
            list: A list of slide content chunks
        """
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Find all headings, paragraphs, lists, tables
        elements = []
        
        # Process headings
        for heading_level in range(1, 7):
            for heading in soup.find_all(f'h{heading_level}'):
                elements.append({
                    'type': 'heading',
                    'level': heading_level,
                    'text': heading.get_text(),
                    'element': heading
                })
        
        # Process paragraphs
        for p in soup.find_all('p'):
            # Skip empty paragraphs
            if not p.get_text().strip():
                continue
                
            text = p.get_text()
            word_count = len(text.split())
            
            # If paragraph is short, keep it as is
            if word_count <= MAX_WORDS_PER_SLIDE:
                elements.append({
                    'type': 'paragraph',
                    'text': text,
                    'word_count': word_count,
                    'element': p
                })
            else:
                # Split long paragraph into sentences
                sentences = sent_tokenize(text)
                current_chunk = []
                current_word_count = 0
                
                for sentence in sentences:
                    sentence_word_count = len(sentence.split())
                    
                    if current_word_count + sentence_word_count <= MAX_WORDS_PER_SLIDE:
                        current_chunk.append(sentence)
                        current_word_count += sentence_word_count
                    else:
                        # Create a new paragraph with current chunk
                        if current_chunk:
                            elements.append({
                                'type': 'paragraph',
                                'text': ' '.join(current_chunk),
                                'word_count': current_word_count,
                                'element': p
                            })
                        
                        # Start a new chunk with this sentence
                        current_chunk = [sentence]
                        current_word_count = sentence_word_count
                
                # Add the last chunk if there is one
                if current_chunk:
                    elements.append({
                        'type': 'paragraph',
                        'text': ' '.join(current_chunk),
                        'word_count': current_word_count,
                        'element': p
                    })
        
        # Process lists
        for list_type in ['ul', 'ol']:
            for list_elem in soup.find_all(list_type):
                list_items = list_elem.find_all('li')
                
                # Process lists that need to be split across slides
                if len(list_items) > MAX_BULLET_POINTS_PER_SLIDE:
                    # Split into chunks of MAX_BULLET_POINTS_PER_SLIDE
                    for i in range(0, len(list_items), MAX_BULLET_POINTS_PER_SLIDE):
                        chunk_items = list_items[i:i + MAX_BULLET_POINTS_PER_SLIDE]
                        elements.append({
                            'type': 'list',
                            'list_type': list_type,
                            'items': [item.get_text() for item in chunk_items],
                            'element': list_elem
                        })
                else:
                    elements.append({
                        'type': 'list',
                        'list_type': list_type,
                        'items': [item.get_text() for item in list_items],
                        'element': list_elem
                    })
        
        # Process tables
        for table in soup.find_all('table'):
            # Extract table header
            header_cells = []
            header_row = table.find('thead')
            if header_row:
                header_cells = [cell.get_text().strip() for cell in header_row.find_all(['th', 'td'])]
            
            # Extract table rows
            rows = []
            for row in table.find_all('tr'):
                # Skip header row
                if row.parent.name == 'thead':
                    continue
                    
                cells = [cell.get_text().strip() for cell in row.find_all(['td', 'th'])]
                if cells:  # Skip empty rows
                    rows.append(cells)
            
            # Split large tables across multiple slides
            if len(rows) > MAX_TABLE_ROWS_PER_SLIDE:
                # Split into chunks
                for i in range(0, len(rows), MAX_TABLE_ROWS_PER_SLIDE):
                    chunk_rows = rows[i:i + MAX_TABLE_ROWS_PER_SLIDE]
                    elements.append({
                        'type': 'table',
                        'header': header_cells,
                        'rows': chunk_rows,
                        'part': f"{i // MAX_TABLE_ROWS_PER_SLIDE + 1} of {(len(rows) + MAX_TABLE_ROWS_PER_SLIDE - 1) // MAX_TABLE_ROWS_PER_SLIDE}",
                        'element': table
                    })
            else:
                elements.append({
                    'type': 'table',
                    'header': header_cells,
                    'rows': rows,
                    'element': table
                })
                
        return elements
    
    def _group_into_slides(self, elements):
        """
        Group analyzed elements into logical slides.
        
        Args:
            elements (list): List of content elements
            
        Returns:
            list: List of slide definitions
        """
        slides = []
        current_slide = {'elements': [], 'word_count': 0}
        current_heading = None
        
        for element in elements:
            # H1 and H2 headings always start a new slide
            if element['type'] == 'heading' and element['level'] <= 2:
                # Save current slide if it has content
                if current_slide['elements']:
                    slides.append(current_slide)
                
                # Start a new slide with this heading
                current_heading = element['text']
                current_slide = {
                    'title': element['text'],
                    'elements': [],
                    'word_count': 0
                }
                continue
            
            # Lower level headings (H3+) might start a new slide if current slide is getting full
            if element['type'] == 'heading' and element['level'] > 2:
                # Check if adding this to current slide would make it too large
                if current_slide['word_count'] > MAX_WORDS_PER_SLIDE / 2:
                    # Save current slide
                    if current_slide['elements']:
                        slides.append(current_slide)
                    
                    # Start a new slide with this heading
                    current_slide = {
                        'title': current_heading or element['text'],
                        'subtitle': element['text'] if current_heading else None,
                        'elements': [],
                        'word_count': 0
                    }
                else:
                    # Add to current slide
                    current_slide['elements'].append(element)
                continue
            
            # Tables and lists always get their own slide
            if element['type'] in ['table', 'list']:
                # Save current slide if it has content
                if current_slide['elements']:
                    slides.append(current_slide)
                
                # Create a slide just for this element
                slides.append({
                    'title': current_heading or "Information",
                    'elements': [element],
                    'word_count': sum(len(item.split()) for item in element.get('items', [])) if element['type'] == 'list' else 0
                })
                
                # Start a new slide for subsequent content
                current_slide = {
                    'title': current_heading,
                    'elements': [],
                    'word_count': 0
                }
                continue
            
            # Paragraphs might need to be split across slides
            if element['type'] == 'paragraph':
                # If this paragraph would make the slide too large, start a new slide
                if current_slide['word_count'] + element['word_count'] > MAX_WORDS_PER_SLIDE:
                    # Save current slide if it has content
                    if current_slide['elements']:
                        slides.append(current_slide)
                    
                    # Start a new slide with this paragraph
                    current_slide = {
                        'title': current_heading or "Information",
                        'elements': [element],
                        'word_count': element['word_count']
                    }
                else:
                    # Add to current slide
                    current_slide['elements'].append(element)
                    current_slide['word_count'] += element['word_count']
        
        # Add the last slide if it has content
        if current_slide['elements']:
            slides.append(current_slide)
            
        return slides
    
    def _create_title_slide(self):
        """Create the title slide for the presentation."""
        slide = self.prs.slides.add_slide(self.layouts['title'])
        
        # Set title
        title = slide.shapes.title
        title.text = f"{self.company_name} Analysis"
        title.text_frame.paragraphs[0].font.size = Pt(54)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.color_scheme['primary']
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        today = datetime.now().strftime("%B %d, %Y")
        subtitle.text = f"Comprehensive Company Analysis\n{today}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.color_scheme['secondary']
        
        # Add a shape for visual interest
        left = Inches(1)
        top = Inches(5.5)
        width = Inches(14)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_scheme['accent']
        shape.line.color.rgb = self.color_scheme['accent']
    
    def _create_section_slide(self, section_title):
        """Create a section divider slide."""
        slide = self.prs.slides.add_slide(self.layouts['section'])
        
        # Set title
        title = slide.shapes.title
        title.text = section_title
        title.text_frame.paragraphs[0].font.size = Pt(54)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.color_scheme['primary']
        
        # Add a visual element
        left = Inches(2)
        top = Inches(4)
        width = Inches(12)
        height = Inches(0.2)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color_scheme['accent']
        shape.line.color.rgb = self.color_scheme['accent']
    
    def _create_content_slide(self, slide_content):
        """Create a content slide based on the provided content."""
        slide = self.prs.slides.add_slide(self.layouts['title_content'])
        
        # Set title
        title = slide.shapes.title
        title.text = slide_content.get('title', '')
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.color_scheme['primary']
        
        # If there's a subtitle, add it
        if slide_content.get('subtitle'):
            subtitle_shape = slide.shapes.add_textbox(
                left=Inches(1),
                top=Inches(1.5),
                width=Inches(14),
                height=Inches(0.6)
            )
            subtitle_frame = subtitle_shape.text_frame
            subtitle_p = subtitle_frame.add_paragraph()
            subtitle_p.text = slide_content['subtitle']
            subtitle_p.font.size = Pt(28)
            subtitle_p.font.color.rgb = self.color_scheme['secondary']
        
        # Process different types of content
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        
        for element in slide_content['elements']:
            if element['type'] == 'paragraph':
                p = tf.add_paragraph()
                p.text = element['text']
                p.font.size = Pt(18)
                p.font.color.rgb = self.color_scheme['text']
                p.space_after = Pt(12)
            
            elif element['type'] == 'heading':
                p = tf.add_paragraph()
                p.text = element['text']
                p.font.size = Pt(24 - element['level'] * 2)  # Size based on heading level
                p.font.bold = True
                p.font.color.rgb = self.color_scheme['primary' if element['level'] <= 2 else 'secondary']
                p.space_before = Pt(12)
                p.space_after = Pt(6)
        
    def _create_list_slide(self, slide_content):
        """Create a slide with a list."""
        slide = self.prs.slides.add_slide(self.layouts['title_content'])
        
        # Set title
        title = slide.shapes.title
        title.text = slide_content.get('title', '')
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.color_scheme['primary']
        
        # Get the list element
        list_element = next((e for e in slide_content['elements'] if e['type'] == 'list'), None)
        if not list_element:
            return
        
        # Process list
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        
        for item_text in list_element['items']:
            p = tf.add_paragraph()
            p.text = item_text
            p.font.size = Pt(18)
            p.font.color.rgb = self.color_scheme['text']
            p.level = 1  # Bullet level
    
    def _create_table_slide(self, slide_content):
        """Create a slide with a table."""
        slide = self.prs.slides.add_slide(self.layouts['title_content'])
        
        # Set title
        title = slide.shapes.title
        title.text = slide_content.get('title', '')
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.color_scheme['primary']
        
        # Get the table element
        table_element = next((e for e in slide_content['elements'] if e['type'] == 'table'), None)
        if not table_element:
            return
        
        # If this is a split table, add part info
        if 'part' in table_element:
            subtitle_shape = slide.shapes.add_textbox(
                left=Inches(11),
                top=Inches(0.9),
                width=Inches(4),
                height=Inches(0.5)
            )
            subtitle_p = subtitle_shape.text_frame.add_paragraph()
            subtitle_p.text = f"Table {table_element['part']}"
            subtitle_p.font.size = Pt(16)
            subtitle_p.font.italic = True
            subtitle_p.font.color.rgb = self.color_scheme['light_text']
            subtitle_p.alignment = PP_ALIGN.RIGHT
        
        # Calculate table size and position
        rows = len(table_element['rows']) + (1 if table_element['header'] else 0)
        cols = max(
            len(table_element['header']),
            max(len(row) for row in table_element['rows']) if table_element['rows'] else 0
        )
        
        # Default content placeholder is too small for most tables, so delete it and add our own
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                sp = shape.element
                sp.getparent().remove(sp)
                break
        
        # Add table with appropriate size
        left = Inches(1)
        top = Inches(2)
        width = Inches(14)
        height = Inches(5.5)
        
        # Create table
        table = slide.shapes.add_table(
            rows, cols,
            left, top, width, height
        ).table
        
        # Set table style
        # Set header
        if table_element['header']:
            for i, text in enumerate(table_element['header']):
                if i < cols:  # Ensure we don't go out of bounds
                    cell = table.cell(0, i)
                    cell.text = text
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.color_scheme['primary']
                    
                    # Style header text
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(16)
                        paragraph.alignment = PP_ALIGN.CENTER
        
        # Set rows
        start_row = 1 if table_element['header'] else 0
        for i, row_data in enumerate(table_element['rows']):
            for j, text in enumerate(row_data):
                if j < cols:  # Ensure we don't go out of bounds
                    cell = table.cell(start_row + i, j)
                    cell.text = text
                    
                    # Style cell text
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.color.rgb = self.color_scheme['text']
                        
                    # Alternate row colors for better readability
                    if i % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    def generate(self):
        """Generate the presentation from markdown content."""
        # Extract content from markdown files
        self._extract_content_from_markdown()
        
        # Create title slide
        self._create_title_slide()
        
        # Create slides for each section in order
        for section_name in self.section_order:
            if section_name in self.sections:
                section = self.sections[section_name]
                
                # Add section divider slide
                self._create_section_slide(section['title'])
                
                # Analyze content and group into slides
                elements = self._analyze_content(section['html'])
                slides = self._group_into_slides(elements)
                
                # Create slides for this section
                for slide_content in slides:
                    # Determine slide type based on content
                    list_element = next((e for e in slide_content['elements'] if e['type'] == 'list'), None)
                    table_element = next((e for e in slide_content['elements'] if e['type'] == 'table'), None)
                    
                    if table_element:
                        self._create_table_slide(slide_content)
                    elif list_element:
                        self._create_list_slide(slide_content)
                    else:
                        self._create_content_slide(slide_content)
    
    def save(self, output_path=None):
        """
        Save the presentation to a file.
        
        Args:
            output_path (str, optional): Path to save the presentation.
                If not provided, saves to the misc directory in the company directory.
        
        Returns:
            str: Path where the presentation was saved
        """
        if not output_path:
            # Create misc directory if it doesn't exist
            misc_dir = os.path.join(self.company_dir, "misc")
            os.makedirs(misc_dir, exist_ok=True)
            
            # Generate filename based on company name and date
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{self.company_name}_Presentation_{timestamp}.pptx"
            output_path = os.path.join(misc_dir, filename)
        
        # Save the presentation
        self.prs.save(output_path)
        return output_path

def generate_presentation(company_dir, output_path=None, color_scheme="professional"):
    """
    Generate a presentation for a company from markdown files.
    
    Args:
        company_dir (str): Directory containing markdown files
        output_path (str, optional): Path to save the presentation
        color_scheme (str, optional): Color scheme name
    
    Returns:
        str: Path where the presentation was saved
    """
    generator = PresentationGenerator(company_dir, color_scheme)
    generator.generate()
    return generator.save(output_path)

def main():
    """Command line interface for presentation generator."""
    parser = argparse.ArgumentParser(description='Generate a presentation from markdown files.')
    parser.add_argument('company_dir', help='Directory containing markdown files')
    parser.add_argument('--output', '-o', help='Output path for the presentation')
    parser.add_argument('--color-scheme', '-c', choices=COLOR_SCHEMES.keys(), default="professional",
                        help='Color scheme for the presentation')
    
    args = parser.parse_args()
    
    try:
        output_path = generate_presentation(args.company_dir, args.output, args.color_scheme)
        print(f"Presentation saved to: {output_path}")
    except Exception as e:
        print(f"Error generating presentation: {e}")
        return 1
    
    return 0

if __name__ == "__main__":
    main() 